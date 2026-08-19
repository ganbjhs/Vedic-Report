"""Documents for TEMPLATE styles — pages designed in Canva (or anything that
exports PNG) with screenshot / text slots drawn on top in the app.

Same inputs as prof_builder (`results`, prepared JPEGs, `layout.placements`),
different drawing: every page starts as the designed image scaled to the paper,
screenshots are fitted (never cropped) into the slots in reading order, and
text slots print report fields. Nothing here reads a screenshot's pixels
beyond scaling — the capture is untouched, this is presentation only.

Fidelity by output (both exact — see registry.OUTPUTS for why these two):
  * PDF  — background image, screenshots and text drawn at their slots.
  * PPTX — the same page as NATIVE OBJECTS: one slide per page, the page art a
           full-bleed picture at the back, then every slot as a picture, a text
           box or a table that can be moved and edited in PowerPoint, Keynote
           or Google Slides. Nothing is flattened into the background.

Two honest limits of the PPTX, both announced on stdout rather than left to be
discovered in a deck (RULEBOOK rule 17):

  * a slide references a typeface by NAME; there is no font file inside a
    .pptx. A style's uploaded font renders as itself on a machine that has it
    installed and is substituted on one that does not — where the PDF, which
    embeds the file, always shows the real face.
  * Helvetica has no file on most machines either; PowerPoint substitutes
    Arial, which is metrically identical, so the slide is named for Arial
    directly instead of pretending.
"""
import datetime
import re
from pathlib import Path

import registry
import layout

# One human name per text field, so the designer, the Canva guide and this
# builder cannot drift apart on what a slot is called.
FIELD_LABELS = {
    "title": "Report title", "date": "Date", "page": "Page no.",
    "pages": "Pages", "index": "#", "account_name": "Account",
    "post_link": "Post URL", "category": "Category", "metrics": "Metrics",
    "handle": "Handle name", "section": "Section", "post_no": "Post 1",
    "post_total": "Top 9 Posts", "post_total_n": "9 Posts",
    "platform": "Platform", "link": "LINK",
    **{f"metric.{k}": f"{k.title()} value" for k in
       ("like", "impressions", "views", "reach", "comments", "shares",
        "followers", "reactions")},
}


def _text_items(profile, page_kind):
    for t in (profile["template"].get("text") or []):
        pg = t.get("page", "post")
        if pg == page_kind or pg == "all":
            yield t


_LOGOS = Path(__file__).resolve().parent / "assets" / "logos"
_PLATFORM_LABEL = {"x": "X", "facebook": "Facebook", "instagram": "Instagram"}


def _field_value(field, ctx):
    if field.startswith("metric."):
        return str((ctx.get("metrics_dict") or {}).get(field[7:], "") or "")
    if field == "metrics":
        m = ctx.get("metrics_dict") or {}
        return " · ".join(f"{k.title()}: {v}" for k, v in m.items() if v and not k.startswith("_")) or ""
    if field == "link":
        return "LINK" if ctx.get("post_link") else ""
    if field == "handle":
        return str(ctx.get("account_name") or "")
    if field == "section":
        return "" if ctx.get("category") in (None, "", "Uncategorized") else str(ctx["category"])
    if field == "platform":
        return _PLATFORM_LABEL.get(ctx.get("platform") or "", "")
    if field == "post_no":
        return f"Post {ctx.get('post_no', '')}"
    if field == "post_total":
        n = ctx.get("post_total", "")
        return f"Top {n} Post{'s' if n != 1 else ''}"
    if field == "post_total_n":
        n = ctx.get("post_total", "")
        return f"{n} Post{'s' if n != 1 else ''}"
    v = ctx.get(field, "")
    return "" if v is None else str(v)


def _text_value(t, ctx):
    """The string a text slot prints: field value, with `label` in front when
    there is a value. Empty → the slot (and its pill) is not drawn at all."""
    if t["field"] == "static":
        return (t.get("label") or "").strip()
    value = _field_value(t["field"], ctx)
    if not value:
        return ""
    if t["field"] == "link" and t.get("label"):
        return t["label"].strip()            # "Open post", not "Open post LINK"
    return f"{t.get('label') or ''}{value}"


def _pill_parts(t, ctx):
    """(label_text, value_text) for a TWO-TONE pill (pill2 + pill). The label
    part is the slot's label; the value part is the field value (or the word
    LINK for a link slot). Both empty → draw nothing."""
    if t["field"] == "static":
        return (t.get("label") or "").strip(), ""
    value = _field_value(t["field"], ctx)
    if not value:
        return "", ""
    return (t.get("label") or "").strip(), value.strip()


_GRID_DEFAULTS = {"box": {"x": 0.02, "y": 0.17, "w": 0.96, "h": 0.80},
                  "cols": 4, "rows": 2, "gap": 0.012, "border": "#222222"}


def _grid_spec(profile):
    g = (profile.get("template") or {}).get("grid")
    if not g:
        return None
    out = dict(_GRID_DEFAULTS)
    out.update({k: v for k, v in g.items() if v is not None})
    out["box"] = {**_GRID_DEFAULTS["box"], **(g.get("box") or {})}
    return out


def _is_grid_section(profile, r):
    g = (profile.get("template") or {}).get("grid")
    if not g:
        return False
    return bool(re.search(g["match"], str(r.get("category") or ""), re.I))


def _grid_cells(spec, W, H):
    """Cell rectangles (x, y_top, w, h) in points for one grid page."""
    box, cols, rows, gap = spec["box"], int(spec["cols"]), int(spec["rows"]), float(spec["gap"])
    bx, by, bw, bh = box["x"] * W, box["y"] * H, box["w"] * W, box["h"] * H
    gx, gy = gap * W, gap * H
    cw = (bw - gx * (cols - 1)) / cols
    ch = (bh - gy * (rows - 1)) / rows
    return [(bx + c * (cw + gx), by + r * (ch + gy), cw, ch)
            for r in range(rows) for c in range(cols)]


def _grid_pages(items, spec):
    """[(section, [(r, img, place), ...]), ...] — one entry per grid PAGE,
    keeping each section on its own page(s)."""
    per = int(spec["cols"]) * int(spec["rows"])
    # group by section (stable: first-seen section order, sheet order inside)
    order = []
    for r, _, _ in items:
        sec = str(r.get("category") or "")
        if sec not in order:
            order.append(sec)
    items = sorted(items, key=lambda x: order.index(str(x[0].get("category") or "")))
    out, cur, bucket = [], None, []
    for r, img, pl in items:
        sec = str(r.get("category") or "")
        if sec != cur or len(bucket) >= per:
            if bucket:
                out.append((cur, bucket))
            cur, bucket = sec, []
        bucket.append((r, img, pl))
    if bucket:
        out.append((cur, bucket))
    return out


def _grid_image(r, composed):
    """The raw screenshot for a grid cell (a comment is short; the cover-cropped
    composite would pad it to the slot's tall shape), else the composite."""
    raw = r.get("screenshot")
    return raw if raw and Path(raw).exists() else composed


def _fit_in(img_path, x, y, w, h):
    """Top-left anchored fit of an image file inside (x, y, w, h) points."""
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    s = min(w / iw, h / ih)
    return x, y, iw * s, ih * s


def _pill_box(t, W, H):
    """(x, y_top, w, h) of a text slot's pill in top-left points."""
    return t["x"] * W, t["y"] * H, t["w"] * W, max(float(t.get("h", 0)) * H, float(t.get("size_pt", 10)) * 1.6)


def _pill_ink(t):
    """Text on a pill: the slot's colour, or white on a dark pill / near-black
    on a light one when the slot did not say."""
    if t.get("color"):
        return t["color"]
    v = (t.get("pill") or "#000000").lstrip("#")
    r, g, b = (int(v[i:i + 2], 16) / 255 for i in (0, 2, 4))
    return "#FFFFFF" if 0.2126 * r + 0.7152 * g + 0.0722 * b < 0.55 else "#111111"


def _metrics_of(r):
    """Sheet columns first (typed by the team from Insights), then whatever a
    capture engine read (influencer). Keys: like, impressions, views, reach…"""
    m = {}
    cap = r.get("metrics") or {}
    for k, v in cap.items():
        if not k.startswith("_") and v not in (None, "", "—"):
            m[k] = v
    if "reactions" in m and "like" not in m:
        m["like"] = m["reactions"]
    for k, v in (r.get("sheet_metrics") or {}).items():
        if v not in (None, ""):
            m[k] = v
    return m


def _sections(results):
    """[(section, count)] in first-seen order, and per-result (post_no, total)."""
    order, counts = [], {}
    for r in results:
        sec = r.get("category") or "Uncategorized"
        if sec not in counts:
            order.append(sec)
            counts[sec] = 0
        counts[sec] += 1
    seen = {}
    per = []
    for r in results:
        sec = r.get("category") or "Uncategorized"
        seen[sec] = seen.get(sec, 0) + 1
        per.append((seen[sec], counts[sec]))
    return [(sec, counts[sec]) for sec in order], per


def _summary_rows(results):
    """The summary table, exactly as both outputs print it."""
    sections, _ = _sections(results)
    return sections + [("Total Items Tracked", len(results))]


# The summary box the PDF falls back to when a style does not place one. Shared
# so the deck cannot end up with the table somewhere the PDF does not have it.
_DEFAULT_SUMMARY_BOX = {"x": 0.36, "y": 0.45, "w": 0.6, "h": 0.5}
_SUMMARY_INK = "#7A3E12"
_SUMMARY_RULE = "#8B5E34"
_SUMMARY_BAND = "#FBEFE0"
_SUMMARY_DIVIDER = 0.62         # the label / value split, as a fraction of width


def _summary_metrics(box, n_rows, page_h_pt):
    """(row height, font size) in points — one rule, two renderers."""
    bh = box["h"] * page_h_pt
    rh = min(bh / max(1, n_rows), 0.42 * 72.0)
    return rh, max(8, min(16, rh * 0.5))


def _trim(value, font, size, max_w_pt):
    """`value` shortened with an ellipsis until it fits `max_w_pt`.

    Measured with reportlab's metrics for BOTH outputs, so a handle that the
    PDF prints as "Kashi Ke Wa…" reads the same on the slide. A slide could
    reflow instead, but then the two documents would disagree about the value,
    which is the one thing they must not do.
    """
    from reportlab.pdfbase.pdfmetrics import stringWidth
    while len(value) > 1 and stringWidth(value, font, size) > max_w_pt:
        value = value[:-2] + "…"
    return value


def _post_ctx(base_ctx, r, page_no, post_no, post_total):
    return {**base_ctx, "page": page_no, "index": post_no, "post_no": post_no,
            "post_total": post_total, "account_name": r.get("account_name", ""),
            "post_link": (r.get("post_link") or r.get("url") or ""),
            "category": r.get("category", ""), "platform": r.get("platform", ""),
            "metrics_dict": _metrics_of(r)}


def _logo_path(platform):
    p = _LOGOS / f"{platform}.png"
    return str(p) if p.exists() else None


def _pages(results, places):
    """[(page_index, [(result, image, place), ...]), ...] in order."""
    out, cur, bucket = [], -1, []
    for r, img, pl in zip(results, places[0], places[1]):
        if pl.page != cur:
            if bucket:
                out.append((cur, bucket))
            cur, bucket = pl.page, []
        bucket.append((r, img, pl))
    if bucket:
        out.append((cur, bucket))
    return out


def _bg(profile, kind):
    p = registry.asset_path(profile, kind)
    return str(p) if p and Path(p).exists() else None


# --------------------------------------------------------------------------- #
# Template fonts — up to 3 uploaded files per style, living beside the page
# images so the runner's copy of assets/<slug>/ carries them into the job.
# --------------------------------------------------------------------------- #
def _font_key(name):
    return "tplf-" + re.sub(r"[^A-Za-z0-9]+", "-", name).strip("-")


def register_fonts(profile) -> dict:
    """{filename: reportlab font name} for every font that really loaded.

    A font that fails to register is REPORTED and dropped, never silently
    swapped — the text still prints in Helvetica, and the log says which file
    the renderer could not read (RULEBOOK rule 17).
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    out = {}
    for name in (profile.get("template") or {}).get("fonts") or []:
        path = registry.font_path(profile, name)
        if not path:
            print(f"[tpl] font {name!r} is not in this style's assets — "
                  "falling back to Helvetica", flush=True)
            continue
        key = _font_key(name)
        try:
            pdfmetrics.registerFont(TTFont(key, str(path)))
            out[name] = key
        except Exception as e:
            print(f"[tpl] could not register font {name!r}: {e} — "
                  "falling back to Helvetica", flush=True)
    return out


def _pdf_font(t, fonts):
    """The reportlab font name for one text slot. An uploaded TTF has one
    weight, so `bold` only picks Helvetica-Bold — a designer who wants bold art
    uploads the bold file and selects it."""
    name = fonts.get(t.get("font") or "")
    if name:
        return name
    return "Helvetica-Bold" if t.get("bold") else "Helvetica"


# --------------------------------------------------------------------------- #
# Non-Latin text — RULEBOOK rule 14: register a Unicode TTF or it prints boxes
# --------------------------------------------------------------------------- #
# This is not hypothetical here. A combined report takes the account name from
# the CAPTURE when the sheet has no handle column, and a Varanasi Facebook Page
# is called "काशी के मोदी". Helvetica is a Type 1 font encoded WinAnsi, so that
# name came out as a row of black rectangles — legible nowhere, and exactly the
# failure rule 14 was written for. The same list `inf_report_builder` uses, plus
# faces that actually carry Devanagari.
_UNICODE_FONTS = (
    "/Library/Fonts/Arial Unicode.ttf",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansDevanagari-Regular.ttf",
    "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
    "/usr/share/fonts/truetype/lohit-devanagari/Lohit-Devanagari.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
    "/Library/Fonts/DejaVuSans.ttf",
)
_UNI_CACHE = []          # [(reportlab font name, {codepoints}) | None]


def _unicode_font():
    """A registered Unicode TTF and the code points it can draw, or None.

    Looked up once per process and never guessed at: the coverage set comes from
    the font's own cmap, so a font that is present but has no Devanagari is
    rejected here instead of printing boxes later.
    """
    if _UNI_CACHE:
        return _UNI_CACHE[0]
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    found = None
    for path in _UNICODE_FONTS:
        p = Path(path)
        if not p.exists():
            continue
        try:
            face = TTFont("tpl-unicode", str(p))
            pdfmetrics.registerFont(face)
            found = ("tpl-unicode", set(face.face.charToGlyph))
            break
        except Exception as e:
            print(f"[tpl] unicode font {p.name} would not load: {e}", flush=True)
    if found is None:
        print("[tpl] no Unicode font on this machine — non-Latin text will "
              "print as boxes (RULEBOOK rule 14)", flush=True)
    _UNI_CACHE.append(found)
    return found


def _fits_font(value, font) -> bool:
    """Can `font` actually draw every character of `value`?"""
    from reportlab.pdfbase import pdfmetrics
    if font.startswith("Helvetica") or font.startswith("Times") or \
            font.startswith("Courier"):
        try:                      # the standard fonts are WinAnsi-encoded
            value.encode("cp1252")
            return True
        except UnicodeEncodeError:
            return False
    try:
        cmap = pdfmetrics.getFont(font).face.charToGlyph
    except Exception:
        return True               # unknown shape — let reportlab decide
    return all(ord(ch) in cmap for ch in value)


def _drawable(value, font):
    """`font`, or a Unicode substitute when `font` cannot draw `value`.

    Falling back silently is what rule 17 forbids, so the swap is announced the
    first time it happens and the text is left alone when nothing can draw it.
    """
    if not value or _fits_font(value, font):
        return font
    uni = _unicode_font()
    if uni and all(ord(ch) in uni[1] for ch in value):
        return uni[0]
    print(f"[tpl] {font} cannot draw {value!r} and no installed font can — "
          "printing it anyway", flush=True)
    return font


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #
def build_pdf(results, images, places, profile, title, out):
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdfcanvas

    pw_in, ph_in = registry.page_inches(profile["page"])
    W, H = pw_in * inch, ph_in * inch
    c = pdfcanvas.Canvas(str(out), pagesize=(W, H))
    c.setTitle(title)
    date = datetime.date.today().strftime("%d-%m-%Y")
    n_pages = 0

    def paint_bg(kind):
        bg = _bg(profile, kind) or _bg(profile, "post")
        if bg:
            c.drawImage(bg, 0, 0, width=W, height=H, mask="auto")

    fonts = register_fonts(profile)

    def draw_text(t, ctx):
        value = _text_value(t, ctx)
        if not value:
            return
        size = float(t.get("size_pt", 10))
        # The value decides the font, not only the slot: a captured Facebook
        # Page name may be Devanagari where the style asked for Helvetica.
        font = _drawable(value, _pdf_font(t, fonts))
        c.setFont(font, size)
        x, y, w = t["x"] * W, H - t["y"] * H - size, t["w"] * W
        if t.get("pill"):
            # a rounded pill behind the text, text centred in it (v3)
            px, py, pw_, ph_ = _pill_box(t, W, H)
            y = H - py - ph_ / 2 - size * 0.35
            if t.get("pill2"):
                # two-tone: label pill (pill2) | value pill (pill), 58/42
                lab, val = _pill_parts(t, ctx)
                lw = pw_ * 0.58 if val else pw_
                c.setFillColor(colors.HexColor(t["pill2"]))
                c.roundRect(px, H - py - ph_, lw, ph_, ph_ / 2, stroke=0, fill=1)
                c.setFillColor(colors.HexColor(_pill_ink({"pill": t["pill2"], "color": t.get("color")})))
                c.drawCentredString(px + lw / 2, y, _trim(lab, font, size, lw - ph_ * 0.6))
                if val:
                    vx = px + lw + ph_ * 0.15
                    vw = pw_ - lw - ph_ * 0.15
                    c.setFillColor(colors.HexColor(t["pill"]))
                    c.roundRect(vx, H - py - ph_, vw, ph_, ph_ / 2, stroke=0, fill=1)
                    c.setFillColor(colors.HexColor(_pill_ink({"pill": t["pill"]})))
                    c.drawCentredString(vx + vw / 2, y, _trim(val, font, size, vw - ph_ * 0.5))
                    if t["field"] in ("post_link", "link") and ctx.get("post_link"):
                        c.linkURL(ctx["post_link"], (vx, H - py - ph_, vx + vw, H - py), relative=0)
                return
            c.setFillColor(colors.HexColor(t["pill"]))
            c.roundRect(px, H - py - ph_, pw_, ph_, ph_ / 2, stroke=0, fill=1)
            c.setFillColor(colors.HexColor(_pill_ink(t)))
            value = _trim(value, font, size, w - ph_ * 0.6)
            c.drawCentredString(x + w / 2, y, value)
            if t["field"] in ("post_link", "link") and ctx.get("post_link"):
                c.linkURL(ctx["post_link"], (px, H - py - ph_, px + pw_, H - py), relative=0)
            return
        c.setFillColor(colors.HexColor(t.get("color") or "#111111"))
        # trim to the slot width — measured in the font it will print in
        value = _trim(value, font, size, w)
        if t.get("align") == "center":
            c.drawCentredString(x + w / 2, y, value)
        elif t.get("align") == "right":
            c.drawRightString(x + w, y, value)
        else:
            c.drawString(x, y, value)
        if t["field"] in ("post_link", "link") and ctx.get("post_link"):
            c.linkURL(ctx["post_link"], (x, y - 2, x + w, y + size), relative=0)

    base_ctx = {"title": title, "date": date, "pages": n_pages}
    tpl = profile["template"]
    # v3: sections that match template.grid.match (counter comments) are laid
    # out MANY per page after the post pages; the rest are one-per-page as
    # before. Numbering (Post i / Top N) counts within the normal ones only.
    gspec = _grid_spec(profile)
    normal = [(r, img, pl) for r, img, pl in zip(results, images, places)
              if not _is_grid_section(profile, r)]
    gridded = [(r, img, pl) for r, img, pl in zip(results, images, places)
               if _is_grid_section(profile, r)]
    pages = _pages([x[0] for x in normal], ([x[1] for x in normal], [x[2] for x in normal]))
    n_pages = len(pages)
    base_ctx["pages"] = n_pages
    sections, per_post = _sections([x[0] for x in normal])

    if _bg(profile, "cover"):
        paint_bg("cover")
        for t in _text_items(profile, "cover"):
            draw_text(t, {**base_ctx, "page": 0})
        c.showPage()

    if _bg(profile, "summary") or tpl.get("summary_box"):
        paint_bg("summary")
        for t in _text_items(profile, "summary"):
            draw_text(t, {**base_ctx, "page": 0})
        box = tpl.get("summary_box") or _DEFAULT_SUMMARY_BOX
        bx, by, bw = box["x"] * W, box["y"] * H, box["w"] * W
        rows = _summary_rows(results)
        rh, fs = _summary_metrics(box, len(rows), H)
        y_top = H - by
        for i, (name, n) in enumerate(rows):
            yy = y_top - (i + 1) * rh
            last = i == len(rows) - 1
            c.setFillColor(colors.HexColor(_SUMMARY_BAND) if i % 2 == 0 else colors.white)
            c.rect(bx, yy, bw, rh, stroke=0, fill=1)
            c.setStrokeColor(colors.HexColor(_SUMMARY_RULE))
            c.setLineWidth(0.6)
            c.rect(bx, yy, bw, rh, stroke=1, fill=0)
            c.line(bx + bw * _SUMMARY_DIVIDER, yy, bx + bw * _SUMMARY_DIVIDER, yy + rh)
            c.setFillColor(colors.HexColor(_SUMMARY_INK))
            label = str(name)[:60]
            c.setFont(_drawable(label, "Helvetica-Bold"), fs)
            c.drawString(bx + 8, yy + rh * 0.32, label)
            c.setFont("Helvetica-Bold" if last else "Helvetica", fs)
            c.setFillColor(colors.HexColor("#111111"))
            c.drawRightString(bx + bw - 10, yy + rh * 0.32, str(n))
        c.showPage()

    idx = 0
    for page_no, (pidx, items) in enumerate(pages, start=1):
        paint_bg("post")
        for r, img, pl in items:
            x = pl.x_in * inch
            y = H - (pl.y_in + pl.h_in) * inch
            c.drawImage(img, x, y, width=pl.w_in * inch, height=pl.h_in * inch,
                        preserveAspectRatio=True, anchor="c", mask="auto")
        # per-post text slots take the FIRST post on the page's fields (single-
        # slot templates, the common case); page-level fields work everywhere.
        r0 = items[0][0]
        post_no, post_total = per_post[idx]
        idx += len(items)
        ctx = _post_ctx(base_ctx, r0, page_no, post_no, post_total)
        for lg in tpl.get("logos") or []:
            logo = _logo_path(r0.get("platform") or "")
            if logo:
                c.drawImage(logo, lg["x"] * W, H - (lg["y"] + lg["h"]) * H,
                            width=lg["w"] * W, height=lg["h"] * H,
                            preserveAspectRatio=True, anchor="c", mask="auto")
        for t in _text_items(profile, "post"):
            draw_text(t, ctx)
        c.showPage()

    # NO trailing links page. A designed report used to end with a sheet-style
    # list of every URL, which was a leftover from the numeric styles: here each
    # post already carries its own LINK, so the list repeated 16 hyperlinks
    # nobody clicked on a page that did not belong to the design. The style's
    # own closing art still prints; `content.links_table` no longer draws
    # anything for a template style (2.4.0).
    # v3: grid pages — counter-comment screenshots, many per page, under the
    # section's name. The "grid" page art if the style has one, else the post
    # art; the heading and anything else comes from text slots on page "grid".
    if gridded and gspec:
        cells = _grid_cells(gspec, W, H)
        for sec, items_ in _grid_pages(gridded, gspec):
            paint_bg("grid")
            gctx = {**base_ctx, "page": n_pages, "category": sec, "section": sec,
                    "post_link": "", "metrics_dict": {}}
            for t in _text_items(profile, "grid"):
                draw_text(t, gctx)
            for (r, img, pl), (cx, cy, cw, ch) in zip(items_, cells):
                img = _grid_image(r, img)
                fx, fy, fw, fh = _fit_in(img, cx, cy, cw, ch)
                c.drawImage(img, fx, H - fy - fh, width=fw, height=fh, mask="auto")
                if gspec.get("border"):
                    c.setStrokeColor(colors.HexColor(gspec["border"]))
                    c.setLineWidth(0.8)
                    c.rect(fx, H - fy - fh, fw, fh, stroke=1, fill=0)
                link = (r.get("post_link") or r.get("url") or "")
                if link:
                    c.linkURL(link, (fx, H - fy - fh, fx + fw, H - fy), relative=0)
            c.showPage()

    if _bg(profile, "end"):
        paint_bg("end")
        for t in _text_items(profile, "end"):
            draw_text(t, {**base_ctx, "page": n_pages})
        c.showPage()
    c.save()


# --------------------------------------------------------------------------- #
# PPTX — the same page, as objects you can still move
# --------------------------------------------------------------------------- #
# A slide names a typeface; it does not carry one. Helvetica has no file on most
# machines, and PowerPoint substitutes Arial for it — metrically the same face,
# which is why the PDF's Helvetica geometry survives the trip. Naming Arial
# outright beats naming a font we know will be swapped.
_PPTX_DEFAULT_FACE = "Arial"
# Arial's ascent in ems. The PDF puts a text slot's BASELINE one em below the
# slot's y; a top-anchored text box at single line spacing puts it one ascent
# below the frame top. Nudging the frame down by the difference makes the two
# documents agree on where the text sits, rather than by a visible hair.
_ASCENT_EM = 0.905


def _rgb(value, default="111111"):
    """'#7A3E12' -> RGBColor. Never raises: a colour the slide cannot parse is
    reported and printed in the default ink rather than taking the deck down."""
    from pptx.dml.color import RGBColor
    raw = str(value or "").lstrip("#").strip()
    if len(raw) == 6:
        try:
            return RGBColor.from_string(raw.upper())
        except ValueError:
            pass
    if raw:
        print(f"[tpl] pptx: {value!r} is not a #RRGGBB colour — using "
              f"#{default}", flush=True)
    return RGBColor.from_string(default)


def face_names(profile) -> dict:
    """{font filename: the family name a slide must ask for}.

    Read from the font FILE, not from its filename: "Brand-Regular.ttf" is
    called "Brand" inside, and a slide asking for "Brand-Regular" would be
    substituted on the designer's own machine.
    """
    from PIL import ImageFont
    out = {}
    for name in (profile.get("template") or {}).get("fonts") or []:
        path = registry.font_path(profile, name)
        if not path:
            print(f"[tpl] pptx: font {name!r} is not in this style's assets — "
                  f"that text uses {_PPTX_DEFAULT_FACE}", flush=True)
            continue
        try:
            family = ImageFont.truetype(str(path), 20).getname()[0]
        except Exception as e:
            print(f"[tpl] pptx: could not read the family name of {name!r}: "
                  f"{e} — that text uses {_PPTX_DEFAULT_FACE}", flush=True)
            continue
        out[name] = family
        print(f"[tpl] pptx: {name!r} is referenced by family name {family!r} — "
              "a machine without it installed will substitute (the PDF embeds "
              "the file itself)", flush=True)
    return out


def _fit_box(img, x_pt, y_pt, w_pt, h_pt):
    """The rectangle reportlab's `preserveAspectRatio` + `anchor='c'` would
    actually paint: fitted inside the box, centred both ways. Same `fit()` the
    placements come from, so the slide lands the picture where the PDF does."""
    import layout
    from PIL import Image
    try:
        with Image.open(img) as im:
            iw, ih = im.size
    except Exception as e:
        print(f"[tpl] pptx: could not measure {img}: {e} — filling the box",
              flush=True)
        return x_pt, y_pt, w_pt, h_pt
    fw, fh = layout.fit(iw, ih, w_pt, h_pt)
    return x_pt + (w_pt - fw) / 2, y_pt + (h_pt - fh) / 2, fw, fh


def _cell_style(cell, fill_hex, line_hex, width_pt=0.6):
    """Solid fill + a rule on all four sides. python-pptx has no border API, so
    the `a:ln*` elements are written directly — they go at the FRONT of tcPr,
    which is where the schema wants them (before the fill)."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(fill_hex)
    tcPr = cell._tc.get_or_add_tcPr()
    emu = int(width_pt * 12700)
    line = _rgb(line_hex)
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):      # reverse: each goes first
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
        tcPr.insert(0, parse_xml(
            f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/'
            f'2006/main" w="{emu}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{line}"/></a:solidFill>'
            f'<a:prstDash val="solid"/></{tag}>'))


def build_pptx(results, images, places, profile, title, out):
    from pptx import Presentation
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    aligns = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
              "right": PP_ALIGN.RIGHT}

    pw_in, ph_in = registry.page_inches(profile["page"])
    W, H = pw_in * 72.0, ph_in * 72.0            # points, exactly like the PDF
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(W), Pt(H)
    blank = prs.slide_layouts[6]                 # the empty layout

    date = datetime.date.today().strftime("%d-%m-%Y")
    fonts = register_fonts(profile)              # reportlab, for measuring only
    faces = face_names(profile)
    tpl = profile["template"]
    gspec = _grid_spec(profile)
    normal = [(r, img, pl) for r, img, pl in zip(results, images, places)
              if not _is_grid_section(profile, r)]
    gridded = [(r, img, pl) for r, img, pl in zip(results, images, places)
               if _is_grid_section(profile, r)]
    pages = _pages([x[0] for x in normal], ([x[1] for x in normal], [x[2] for x in normal]))
    n_pages = len(pages)

    def new_slide(kind):
        """A slide with its page art already at the back — the background is the
        first shape on it, so everything added afterwards sits on top of it."""
        slide = prs.slides.add_slide(blank)
        bg = _bg(profile, kind) or _bg(profile, "post")
        if bg:
            slide.shapes.add_picture(bg, 0, 0, width=Pt(W), height=Pt(H))
        return slide

    def add_picture(slide, img, x_pt, y_pt, w_pt, h_pt):
        x, y, w, h = _fit_box(img, x_pt, y_pt, w_pt, h_pt)
        return slide.shapes.add_picture(img, Pt(x), Pt(y), Pt(w), Pt(h))

    def draw_text(slide, t, ctx):
        value = _text_value(t, ctx)
        if not value:
            return
        size = float(t.get("size_pt", 10))
        # Measured in the font the PDF would print it in, so both documents
        # trim the same value to the same characters. The SLIDE keeps the text
        # as text, so a non-Latin string needs no font swap here — PowerPoint
        # falls back per glyph where reportlab's WinAnsi Helvetica cannot.
        if t.get("pill"):
            # the pill is a real rounded rectangle carrying the text — one
            # object to move, recolour or retype in PowerPoint (v3)
            from pptx.enum.shapes import MSO_SHAPE
            px, py, pw_, ph_ = _pill_box(t, W, H)

            def pill_shape(x0, w0, fill, text_, ink, link=None):
                shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x0), Pt(py), Pt(w0), Pt(ph_))
                shp.adjustments[0] = 0.5
                shp.fill.solid()
                shp.fill.fore_color.rgb = _rgb(fill)
                shp.line.fill.background()
                shp.shadow.inherit = False
                tf = shp.text_frame
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = tf.margin_right = Pt(ph_ * 0.3)
                tf.margin_top = tf.margin_bottom = 0
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = _trim(text_, _drawable(text_, _pdf_font(t, fonts)), size, w0 - ph_ * 0.6)
                run.font.name = faces.get(t.get("font") or "") or _PPTX_DEFAULT_FACE
                run.font.size = Pt(size)
                run.font.bold = bool(t.get("bold"))
                run.font.color.rgb = _rgb(ink)
                if link:
                    run.hyperlink.address = link
                    run.font.underline = False        # the pill IS the button
                return shp

            link = ctx.get("post_link") if t["field"] in ("post_link", "link") else None
            if t.get("pill2"):
                lab, val = _pill_parts(t, ctx)
                lw = pw_ * 0.58 if val else pw_
                pill_shape(px, lw, t["pill2"], lab, _pill_ink({"pill": t["pill2"], "color": t.get("color")}))
                if val:
                    vx = px + lw + ph_ * 0.15
                    pill_shape(vx, pw_ - lw - ph_ * 0.15, t["pill"], val, _pill_ink({"pill": t["pill"]}), link)
                return
            return pill_shape(px, pw_, t["pill"], value, _pill_ink(t), link)
        value = _trim(value, _drawable(value, _pdf_font(t, fonts)), size,
                      t["w"] * W)
        box = slide.shapes.add_textbox(Pt(t["x"] * W),
                                       Pt(t["y"] * H + (1 - _ASCENT_EM) * size),
                                       Pt(t["w"] * W), Pt(size * 1.5))
        tf = box.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = aligns.get(t.get("align") or "left", PP_ALIGN.LEFT)
        para.line_spacing = 1.0
        run = para.add_run()
        run.text = value
        run.font.name = faces.get(t.get("font") or "") or _PPTX_DEFAULT_FACE
        run.font.size = Pt(size)
        run.font.bold = bool(t.get("bold"))
        run.font.color.rgb = _rgb(t.get("color"))
        if t["field"] in ("post_link", "link") and ctx.get("post_link"):
            # A real hyperlink on the run, not a picture of one: the LINK button
            # is clickable in PowerPoint and survives an upload to Slides.
            run.hyperlink.address = ctx["post_link"]
        return box

    base_ctx = {"title": title, "date": date, "pages": n_pages}

    if _bg(profile, "cover"):
        slide = new_slide("cover")
        for t in _text_items(profile, "cover"):
            draw_text(slide, t, {**base_ctx, "page": 0})

    if _bg(profile, "summary") or tpl.get("summary_box"):
        slide = new_slide("summary")
        for t in _text_items(profile, "summary"):
            draw_text(slide, t, {**base_ctx, "page": 0})
        box = tpl.get("summary_box") or _DEFAULT_SUMMARY_BOX
        rows = _summary_rows(results)
        rh, fs = _summary_metrics(box, len(rows), H)
        bx, by, bw = box["x"] * W, box["y"] * H, box["w"] * W
        table = slide.shapes.add_table(len(rows), 2, Pt(bx), Pt(by), Pt(bw),
                                       Pt(rh * len(rows))).table
        # A real table, editable row by row — and its own look, not the theme's
        # banded default, which would repaint the PDF's colours blue.
        table.first_row = table.horz_banding = False
        table.columns[0].width = Pt(bw * _SUMMARY_DIVIDER)
        table.columns[1].width = Pt(bw * (1 - _SUMMARY_DIVIDER))
        for i, (name, n) in enumerate(rows):
            last = i == len(rows) - 1
            table.rows[i].height = Pt(rh)
            cells = ((str(name)[:60], PP_ALIGN.LEFT, _SUMMARY_INK, True),
                     (str(n), PP_ALIGN.RIGHT, "#111111", last))
            for j, (value, align, ink, bold) in enumerate(cells):
                cell = table.cell(i, j)
                _cell_style(cell, _SUMMARY_BAND if i % 2 == 0 else "#FFFFFF",
                            _SUMMARY_RULE)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = Pt(8)
                cell.margin_top = cell.margin_bottom = 0
                para = cell.text_frame.paragraphs[0]
                para.alignment = align
                run = para.add_run()
                run.text = value
                run.font.size = Pt(fs)
                run.font.bold = bold
                run.font.name = _PPTX_DEFAULT_FACE
                run.font.color.rgb = _rgb(ink)

    _sec, per_post = _sections(results)
    idx = 0
    for page_no, (_pidx, items) in enumerate(pages, start=1):
        slide = new_slide("post")
        for r, img, pl in items:
            add_picture(slide, img, pl.x_in * 72.0, pl.y_in * 72.0,
                        pl.w_in * 72.0, pl.h_in * 72.0)
        r0 = items[0][0]
        post_no, post_total = per_post[idx]
        idx += len(items)
        ctx = _post_ctx(base_ctx, r0, page_no, post_no, post_total)
        for lg in tpl.get("logos") or []:
            logo = _logo_path(r0.get("platform") or "")
            if logo:
                add_picture(slide, logo, lg["x"] * W, lg["y"] * H,
                            lg["w"] * W, lg["h"] * H)
        for t in _text_items(profile, "post"):
            draw_text(slide, t, ctx)

    # v3: grid pages (counter comments), same cells as the PDF
    if gridded and gspec:
        cells = _grid_cells(gspec, W, H)
        for sec, items_ in _grid_pages(gridded, gspec):
            slide = new_slide("grid")
            gctx = {**base_ctx, "page": n_pages, "category": sec, "section": sec,
                    "post_link": "", "metrics_dict": {}}
            for t in _text_items(profile, "grid"):
                draw_text(slide, t, gctx)
            for (r, img, pl), (cx, cy, cw, ch) in zip(items_, cells):
                img = _grid_image(r, img)
                fx, fy, fw, fh = _fit_in(img, cx, cy, cw, ch)
                pic = slide.shapes.add_picture(img, Pt(fx), Pt(fy), Pt(fw), Pt(fh))
                if gspec.get("border"):
                    pic.line.color.rgb = _rgb(gspec["border"])
                    pic.line.width = Pt(0.8)
                link = (r.get("post_link") or r.get("url") or "")
                if link:
                    pic.click_action.hyperlink.address = link

    # Same ending as the PDF: the style's closing art if it has one, and no
    # trailing links list (2.4.0 — every post carries its own LINK).
    if _bg(profile, "end"):
        slide = new_slide("end")
        for t in _text_items(profile, "end"):
            draw_text(slide, t, {**base_ctx, "page": n_pages})

    prs.save(str(out))


BUILDERS = {"pdf": build_pdf, "pptx": build_pptx}
