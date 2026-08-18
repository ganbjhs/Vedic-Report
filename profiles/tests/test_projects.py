"""v3: projects + page background + numeric PPTX — zero captures.

Runs against a throwaway DATA_DIR so the developer's real jobs.db and user
styles are never touched.
"""
import json
import os
import shutil
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[1]
TMP = Path(tempfile.mkdtemp(prefix="vr-test-"))
os.environ["DATA_DIR"] = str(TMP / "data")
os.environ.setdefault("APP_USERS", "t:tttttttt")
os.environ.setdefault("SESSION_SECRET", "x" * 40)
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "profiles"))

import registry                                     # noqa: E402
import prof_builder                                 # noqa: E402
from webapp import config, styles as wstyles        # noqa: E402
from webapp.jobs import store                       # noqa: E402

FAILS = []


def check(cond, msg):
    print(("  PASS  " if cond else "  FAIL  ") + msg)
    if not cond:
        FAILS.append(msg)


def main():
    config.ensure_dirs()
    store.init()

    print("\n1. projects: Unsorted exists, CRUD, styles list")
    uns = store.project_by_slug(store.UNSORTED_SLUG)
    check(uns is not None, "Unsorted project is created on init")
    pid = store.project_create("t", "kashi", "Kashi", client="KKW", emoji="🪔")
    check(store.project_get(pid)["name"] == "Kashi", "create + get")
    store.project_set_styles(pid, [{"slug": "client-deck", "outputs": ["pdf", "pptx"]},
                                   {"slug": "combined-16x9", "outputs": []}])
    got = store.project_styles(pid)
    check([g["slug"] for g in got] == ["client-deck", "combined-16x9"], "styles kept in order")
    check(got[0]["outputs"] == ["pdf", "pptx"], "per-style outputs kept")
    jid = store.create("t", "r", "R", "client-deck", 1, "x", project_id=pid)
    check(store.get(jid)["project_id"] == pid, "job carries project_id")
    check(store.list_for_project(pid)[0]["id"] == jid, "list_for_project")
    check(store.project_delete(pid) is False, "a project with runs is not deleted")
    store.init()
    old = store.create("t", "old", "Old", "twitter", 1, "x")   # project_id ''
    store.init()
    check(store.get(old)["project_id"] == uns["id"], "jobs without a project migrate to Unsorted")

    print("\n2. registry: page.background")
    p = registry.load("client-deck")
    p["page"]["background"] = {"color": "#1B2440"}
    registry.validate(p)
    check(True, "colour background validates")
    check(registry.background_color(p) == "#1B2440", "background_color reads it")
    for bad in ({"color": "navy"}, {"image": "../x.png"}, {"nope": 1}, "red"):
        p["page"]["background"] = bad
        try:
            registry.validate(p)
            check(False, f"{bad!r} refused")
        except registry.ProfileError:
            check(True, f"{bad!r} refused")
    p["page"]["background"] = None
    check("pptx" in registry.NUMERIC_OUTPUTS, "numeric styles may build pptx")
    check(registry.load("client-deck")["outputs"] == ["pdf", "docx", "pptx"], "client-deck declares pptx")

    print("\n3. builder: numeric PPTX from the fixture, with a dark background")
    fixture = ROOT / "data" / "sample-fixture"
    res = json.loads((fixture / "results.json").read_text())
    for r in res:
        r["screenshot"] = str(fixture / "screenshots" / Path(r["screenshot"]).name)
    res = [r for r in res if prof_builder.usable(r)]
    check(len(res) >= 3, f"{len(res)} usable fixture posts")
    p = registry.load("client-deck")
    p["page"]["background"] = {"color": "#0F172A"}
    check(prof_builder._palette(p)["ink"] == "#FFFFFF", "dark background → light ink")
    with tempfile.TemporaryDirectory() as td:
        images, places = prof_builder.prepare(res, p, td)
        out = Path(td) / "t.pptx"
        prof_builder.build_pptx(res, images, places, p, "T", out)
        check(out.stat().st_size > 10_000, "pptx written")
        from pptx import Presentation
        prs = Presentation(str(out))
        n_pages = max(pl.page for pl in places) + 1
        expect = n_pages + (1 if p["content"].get("cover") else 0) + (1 if p["content"].get("links_table") else 0)
        check(len(prs.slides) == expect, f"one slide per page (+cover, +links): {len(prs.slides)} == {expect}")
        pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
        check(pics == len(res), f"every post is a picture on a slide ({pics})")
        pdf = Path(td) / "t.pdf"
        prof_builder.build_pdf(res, images, places, p, "T", pdf)
        check(pdf.stat().st_size > 10_000, "pdf with background written")

    print("\n4. styles: set_background + fork_for_project")
    raw = {"schema": 1, "slug": "t-bg", "label": "T bg", "extends": "client-deck",
           "capture": {}, "image": {}, "page": {}, "content": {}, "outputs": ["pdf", "pptx"]}
    wstyles.save(raw)
    wstyles.set_background("t-bg", color="#112233")
    check(json.loads(wstyles._path("t-bg").read_text())["page"]["background"] == {"color": "#112233"}, "colour stored")
    from PIL import Image
    import io
    im = Image.new("RGB", (400, 600), (10, 20, 30)); b = io.BytesIO(); im.save(b, "PNG")
    wstyles.set_background("t-bg", image=b.getvalue())
    check((wstyles.asset_dir("t-bg") / "background.png").exists(), "image stored beside the profile")
    check(registry.background_path(registry.load("t-bg")).exists(), "registry finds it")
    wstyles.set_background("t-bg", remove=True)
    check(not (wstyles.asset_dir("t-bg") / "background.png").exists(), "remove clears the file")
    try:
        wstyles.set_background("client-deck", color="#000000")
        check(False, "a shipped style refuses a background")
    except wstyles.StyleError:
        check(True, "a shipped style refuses a background")
    slug = wstyles.fork_for_project("client-deck", {"slug": "kashi", "name": "Kashi", "id": "abc"})
    check(slug == "client-deck-kashi", f"fork slug {slug}")
    check(wstyles.fork_for_project("client-deck", {"slug": "kashi", "name": "Kashi", "id": "abc"}) == slug, "fork is idempotent")
    check(json.loads(wstyles._path(slug).read_text())["extends"] == "client-deck", "fork extends the source")

    shutil.rmtree(TMP, ignore_errors=True)
    if FAILS:
        print(f"\nFAILED: {len(FAILS)}")
        sys.exit(1)
    print("\nPROJECTS OK — store, background, numeric PPTX")


if __name__ == "__main__":
    main()
