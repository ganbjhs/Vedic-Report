"""Build a profile's documents from results.json — PDF, DOCX and/or PPTX.

Geometry comes from `layout.py`, shape from `shapes.py`; this module only draws.
That split is what makes the whole layout testable with no browser and no
rasteriser (`tests/test_parity.py`).

`src/report_builder.py` and `influencer/inf_report_builder.py` are NOT imported
and NOT modified. The two existing report types keep their own builders and
entrypoints (docs/profile-engine.md §10 decision 3); this one serves additional
profiles.

Usage: python profiles/prof_builder.py "<header>" "<stem>" "<profile-slug>" [pdf,docx,pptx]
"""
import json
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[0].parent if HERE.name != "profiles" else HERE.parent
sys.path.insert(0, str(HERE))
import layout        # noqa: E402
import progress      # noqa: E402
import registry      # noqa: E402
import shapes        # noqa: E402

OUT = ROOT / "reports"

_JPEG_QUALITY = 88          # matches the frozen builders — no visible loss
INCH = 72.0


# --------------------------------------------------------------------------- #
# Inputs
# --------------------------------------------------------------------------- #
def usable(r) -> bool:
    """Only a cleanly captured link belongs in a document.

    Identical rule to the frozen builders: status must be exactly "ok", so every
    demoted status (overlay_blocked, parent_lost, age_restricted, ...) is left
    out with no per-status logic here.
    """
    if r.get("status") != "ok":
        return False
    shot = r.get("screenshot")
    return bool(shot) and Path(shot).exists() and Path(shot).stat().st_size > 0


def _resolve(path_str):
    """results.json holds ABSOLUTE paths, which is correct and load-bearing
    (RULEBOOK rule 2 — the code is copied into the job dir, so they resolve
    inside it). A consumer running here, inside the job, may trust them; we only
    fall back to the local screenshots folder if the file has moved."""
    p = Path(path_str)
    if p.exists():
        return p
    local = OUT / "screenshots" / p.name
    return local if local.exists() else p


def prepare(results, profile, workdir):
    """Composite each master per the profile, then hand back placements.

    Two-pass, as docs/profile-engine.md requires: `bordered`/`shadowed` grow the
    image and therefore change its aspect, so the placement is computed from the
    master, used to scale point-valued decoration, and then recomputed from the
    composed image.
    """
    from PIL import Image

    spec = profile["image"]
    masters = []
    for r in results:
        with Image.open(_resolve(r["screenshot"])) as im:
            masters.append(im.convert("RGB"))

    # pass 1 — provisional placement, only to learn each image's width in inches
    provisional = layout.placements(profile, [m.size for m in masters])

    prepared, dims = [], []
    for i, (im, place) in enumerate(zip(masters, provisional), 1):
        composed = shapes.compose(im, spec, placement_w_in=place.w_in)
        dst = Path(workdir) / f"{i:03d}.jpg"
        composed.convert("RGB").save(dst, "JPEG", quality=_JPEG_QUALITY,
                                     optimize=True, subsampling=0)
        prepared.append(str(dst))
        dims.append(composed.size)
        if composed is not im:
            composed.close()
        im.close()

    # pass 2 — the placement the document actually uses
    return prepared, layout.placements(profile, dims)


def _background(profile):
    """(image_path or None, '#RRGGBB' or None) — the page background, if any.
    A missing image file falls back to the colour (or nothing) and says so,
    rather than taking the build down (rule 17)."""
    img = registry.background_path(profile)
    color = registry.background_color(profile)
    if img is not None and not Path(img).exists():
        print(f"[report] background image {img} is missing — "
              f"{'using the colour' if color else 'printing without it'}", flush=True)
        img = None
    return (str(img) if img else None), color


def _palette(profile):
    """Ink colours for text: dark on a light page, light on a dark one. Decided
    from page.background.color; an image background keeps the dark ink (its
    tone is not known) — set a colour too if the art is dark."""
    bg = registry.background_color(profile)
    if bg:
        r, g, b = _rgb_hex(bg)
        lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
        if lum < 0.45:
            return {"ink": "#FFFFFF", "grey": "#E2E8F0", "faint": "#CBD5E1",
                    "accent": "#7CC4FA"}
    return {"ink": "#0F172A", "grey": "#334155", "faint": "#64748B",
            "accent": "#1D9BF0"}


def _rgb_hex(value, default=(1.0, 1.0, 1.0)):
    v = (value or "").lstrip("#")
    if len(v) != 6:
        return default
    return tuple(int(v[i:i + 2], 16) / 255.0 for i in (0, 2, 4))


def shown_link(r) -> str:
    link = (r.get("post_link") or r.get("url") or "").strip()
    return "" if link.startswith("file://") else link


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #
def build_pdf(results, images, places, profile, title, out):
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdfcanvas

    pw_in, ph_in = registry.page_inches(profile["page"])
    page_w, page_h = pw_in * inch, ph_in * inch
    top, right, bottom, left = profile["page"]["margins_in"]
    content = profile["content"]
    pal = _palette(profile)
    accent = colors.HexColor(pal["accent"])
    ink, grey, faint = (colors.HexColor(pal["ink"]), colors.HexColor(pal["grey"]),
                        colors.HexColor(pal["faint"]))

    c = pdfcanvas.Canvas(str(out), pagesize=(page_w, page_h))
    c.setTitle(title)

    bg_img, bg_color = _background(profile)

    def paint_bg():
        """Called at the START of every page: the fill first, then the image
        stretched to the page (a background is designed at the page's aspect,
        so no fitting — exactly what the PPTX slide does)."""
        if bg_color:
            c.setFillColorRGB(*_rgb_hex(bg_color))
            c.rect(0, 0, page_w, page_h, stroke=0, fill=1)
        if bg_img:
            c.drawImage(bg_img, 0, 0, width=page_w, height=page_h,
                        preserveAspectRatio=False, mask="auto")

    def footer(page_no, pages):
        if not content.get("footer"):
            return
        text = content["footer"].format(page=page_no, pages=pages,
                                        title=title)
        c.setFont("Helvetica", 8)
        c.setFillColor(faint)
        c.drawCentredString(page_w / 2, bottom * inch / 2, text)

    n_pages = (max((p.page for p in places), default=-1) + 1) if places else 0

    if content.get("cover"):
        paint_bg()
        c.setFont("Helvetica-Bold", 26)
        c.setFillColor(ink)
        c.drawCentredString(page_w / 2, page_h * 0.56, title)
        c.setStrokeColor(accent)
        c.setLineWidth(2)
        c.line(page_w * 0.3, page_h * 0.53, page_w * 0.7, page_h * 0.53)
        c.setFont("Helvetica", 11)
        c.setFillColor(faint)
        c.drawCentredString(page_w / 2, page_h * 0.49,
                            f"{len(results)} post(s)")
        c.showPage()

    current = -1
    cursor_y = None          # lowest ink on the current page, for the links table
    for r, img, place in zip(results, images, places):
        if place.page != current:
            if current >= 0:
                footer(current + 1, n_pages)
                c.showPage()
            current = place.page
            cursor_y = None
            paint_bg()
            if content.get("header"):
                c.setFont("Helvetica-Bold", 13)
                c.setFillColor(ink)
                c.drawCentredString(page_w / 2,
                                    page_h - top * inch * 0.55,
                                    content["header"].format(title=title))

        # reportlab's origin is bottom-left; layout works from the top.
        x = place.x_in * inch
        y = page_h - (place.y_in + place.h_in) * inch
        c.drawImage(img, x, y, width=place.w_in * inch, height=place.h_in * inch,
                    preserveAspectRatio=True, anchor="n", mask="auto")

        caption_y = y - 11
        for field in content.get("per_post_fields") or []:
            value = r.get(field) or ""
            if field == "post_link":
                value = shown_link(r)
            if not value:
                continue
            c.setFont("Helvetica", 7.5)
            c.setFillColor(accent if field == "post_link"
                           else grey)
            c.drawString(x, caption_y, str(value)[:110])
            if field == "post_link":
                c.linkURL(value, (x, caption_y - 2, x + place.w_in * inch,
                                  caption_y + 8), relative=0)
            caption_y -= 10

        for label, key in content.get("metrics") or []:
            val = (r.get("metrics") or {}).get(key, "—")
            c.setFont("Helvetica", 7.5)
            c.setFillColor(grey)
            c.drawString(x, caption_y, f"{label}: {val}")
            caption_y -= 9.5

        cursor_y = caption_y if cursor_y is None else min(cursor_y, caption_y)

    if places:
        footer(current + 1, n_pages)

    # The links table FLOWS below the last screenshot rather than starting its
    # own page — matching the frozen builders, which say so explicitly ("no page
    # break before it"). Forcing a page here made the twitter profile 9 pages
    # against the frozen builder's 8, which the geometry parity test could not
    # see because it only compares image placements.
    if content.get("links_table"):
        need = 3 * 12 + 26                     # heading + a few rows
        if cursor_y is None or cursor_y - 26 < bottom * inch + need:
            if places:
                c.showPage()
            paint_bg()
            cursor_y = page_h - top * inch + 14
        y = cursor_y - 26
        c.setFont("Helvetica-Bold", 14)
        c.setFillColor(ink)
        c.drawString(left * inch, y, "Links")
        y -= 18
        c.setFont("Helvetica", 8)
        for r in results:
            link = shown_link(r)
            if not link:
                continue
            if y < bottom * inch:
                c.showPage()
                paint_bg()
                y = page_h - top * inch
                c.setFont("Helvetica", 8)
            c.setFillColor(accent)
            c.drawString(left * inch, y, link[:150])
            c.linkURL(link, (left * inch, y - 2,
                             page_w - right * inch, y + 8), relative=0)
            y -= 12
        c.showPage()
    elif places:
        c.showPage()

    c.save()


# --------------------------------------------------------------------------- #
# DOCX
# --------------------------------------------------------------------------- #
def build_docx(results, images, places, profile, title, out):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    content = profile["content"]
    pw_in, ph_in = registry.page_inches(profile["page"])
    top, right, bottom, left = profile["page"]["margins_in"]

    doc = Document()
    section = doc.sections[0]
    section.page_width, section.page_height = Inches(pw_in), Inches(ph_in)
    section.top_margin, section.bottom_margin = Inches(top), Inches(bottom)
    section.left_margin, section.right_margin = Inches(left), Inches(right)

    if content.get("header"):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(content["header"].format(title=title))
        run.bold = True
        run.font.size = Pt(16)

    per_page = registry.per_page(profile)
    for i, (r, img, place) in enumerate(zip(results, images, places)):
        if i and i % per_page == 0:
            doc.add_page_break()
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(img, width=Inches(place.w_in),
                                height=Inches(place.h_in))
        for field in content.get("per_post_fields") or []:
            value = shown_link(r) if field == "post_link" else (r.get(field) or "")
            if not value:
                continue
            cp = doc.add_paragraph()
            run = cp.add_run(str(value))
            run.font.size = Pt(8)
            if field == "post_link":
                run.font.color.rgb = RGBColor(0x1D, 0x9B, 0xF0)
        for label, key in content.get("metrics") or []:
            mp = doc.add_paragraph()
            mrun = mp.add_run(f"{label}: {(r.get('metrics') or {}).get(key, '—')}")
            mrun.font.size = Pt(8)

    if content.get("links_table"):
        doc.add_heading("Links", level=1)
        table = doc.add_table(rows=1, cols=1)
        table.style = "Table Grid"
        hdr = table.rows[0].cells[0].paragraphs[0].add_run("Link")
        hdr.bold = True
        hdr.font.color.rgb = RGBColor(0x1D, 0x9B, 0xF0)
        for r in results:
            link = shown_link(r)
            table.add_row().cells[0].text = link or "—"

    doc.save(str(out))


# --------------------------------------------------------------------------- #
# PPTX (v3) — the same pages as the PDF, every element a native slide object
# --------------------------------------------------------------------------- #
def build_pptx(results, images, places, profile, title, out):
    """One slide per page, sized exactly like the PDF page (points → EMU).

    Order on every slide: background fill/picture FIRST (so it is at the back),
    then screenshots as pictures at their placements, then caption lines as
    text boxes (the post link as a real hyperlink), then metrics, then the
    header. Cover and Links pages become slides too, so a deck reads like the
    PDF page for page. python-pptx lives in requirements-web.txt; a bare CLI
    install gets a clear message from main() and still gets its PDF.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    pw_in, ph_in = registry.page_inches(profile["page"])
    W, H = pw_in * INCH, ph_in * INCH
    top, right, bottom, left = profile["page"]["margins_in"]
    content = profile["content"]
    pal = _palette(profile)
    _c = lambda h: RGBColor(*(int(round(v * 255)) for v in _rgb_hex(h)))
    accent, ink, grey, faint = (_c(pal["accent"]), _c(pal["ink"]),
                                _c(pal["grey"]), _c(pal["faint"]))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(W), Pt(H)
    blank = prs.slide_layouts[6]
    bg_img, bg_color = _background(profile)

    def new_slide():
        slide = prs.slides.add_slide(blank)
        if bg_color:
            r, g, b = (int(round(v * 255)) for v in _rgb_hex(bg_color))
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(r, g, b)
        if bg_img:
            slide.shapes.add_picture(bg_img, 0, 0, width=Pt(W), height=Pt(H))
        return slide

    def text(slide, value, x, y, w, h, size, color, bold=False,
             align=PP_ALIGN.LEFT, link=None):
        box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = box.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = str(value)
        run.font.size = Pt(size)
        run.font.bold = bool(bold)
        run.font.color.rgb = color
        run.font.name = "Arial"          # Helvetica's metric twin — see RULEBOOK 18a
        if link:
            run.hyperlink.address = link
        return box

    n_pages = (max((p.page for p in places), default=-1) + 1) if places else 0

    def footer(slide, page_no):
        if not content.get("footer"):
            return
        value = content["footer"].format(page=page_no, pages=n_pages, title=title)
        text(slide, value, 0, H - bottom * INCH / 2 - 5, W, 12, 8, faint,
             align=PP_ALIGN.CENTER)

    if content.get("cover"):
        slide = new_slide()
        text(slide, title, 0, H * 0.44 - 18, W, 36, 26, ink, bold=True,
             align=PP_ALIGN.CENTER)
        text(slide, f"{len(results)} post(s)", 0, H * 0.51, W, 16, 11, faint,
             align=PP_ALIGN.CENTER)

    slide, current = None, -1
    for r, img, place in zip(results, images, places):
        if place.page != current:
            if slide is not None:
                footer(slide, current + 1)
            current = place.page
            slide = new_slide()
            if content.get("header"):
                text(slide, content["header"].format(title=title),
                     0, top * INCH * 0.55 - 9, W, 18, 13, ink, bold=True,
                     align=PP_ALIGN.CENTER)
        x, y = place.x_in * INCH, place.y_in * INCH
        w, h = place.w_in * INCH, place.h_in * INCH
        slide.shapes.add_picture(img, Pt(x), Pt(y), Pt(w), Pt(h))
        cap_y = y + h + 3
        for field in content.get("per_post_fields") or []:
            value = shown_link(r) if field == "post_link" else (r.get(field) or "")
            if not value:
                continue
            is_link = field == "post_link"
            text(slide, str(value)[:110], x, cap_y, w, 11, 7.5,
                 accent if is_link else grey, link=(value if is_link else None))
            cap_y += 10
        for label, key in content.get("metrics") or []:
            val = (r.get("metrics") or {}).get(key, "—")
            text(slide, f"{label}: {val}", x, cap_y, w, 11, 7.5, grey)
            cap_y += 9.5
    if slide is not None:
        footer(slide, current + 1)

    if content.get("links_table"):
        links = [shown_link(r) for r in results if shown_link(r)]
        per_slide = max(1, int((H - (top + bottom) * INCH - 30) // 12))
        for start in range(0, max(len(links), 1), per_slide):
            slide = new_slide()
            y = top * INCH
            text(slide, "Links", left * INCH, y, W - (left + right) * INCH, 18,
                 14, ink, bold=True)
            y += 22
            for link in links[start:start + per_slide]:
                text(slide, link[:150], left * INCH, y, W - (left + right) * INCH,
                     11, 8, accent, link=link)
                y += 12
            if not links:
                break

    prs.save(str(out))


BUILDERS = {"pdf": build_pdf, "docx": build_docx, "pptx": build_pptx}


def _mb(path):
    return round(Path(path).stat().st_size / 1_048_576, 1)


def wanted_outputs(profile, asked: str = "") -> list:
    """The profile's outputs, narrowed to the ones this run asked for.

    The style still owns WHICH documents exist — a request is a filter over
    that list and never an addition to it, so a hand-crafted job cannot ask a
    numeric style for a deck. An empty or fully-unrecognised request means
    "everything the style declares", which is what every caller before 2.4.0
    meant by passing nothing.
    """
    declared = list(profile["outputs"])
    want = [w for w in (a.strip().lower() for a in (asked or "").split(","))
            if w]
    if not want:
        return declared
    kept = [k for k in declared if k in want]
    unknown = [w for w in want if w not in declared]
    if unknown:
        print(f"[report] ignoring requested output(s) {unknown} — "
              f"{profile['slug']} builds {declared}", flush=True)
    if not kept:
        print("[report] none of the requested outputs belong to this style — "
              f"building {declared}", flush=True)
        return declared
    return kept


def main():
    title = sys.argv[1] if len(sys.argv) > 1 else "Report"
    stem = sys.argv[2] if len(sys.argv) > 2 else "Report"
    slug = sys.argv[3] if len(sys.argv) > 3 else "twitter"
    asked = sys.argv[4] if len(sys.argv) > 4 else ""
    profile = registry.load(slug)

    all_results = json.loads((OUT / "results.json").read_text())
    results = [r for r in all_results if usable(r)]
    skipped = len(all_results) - len(results)
    if not results:
        print(f"[report] no capturable links ({skipped} skipped) — "
              "nothing to build", flush=True)
        return

    OUT.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as td:
        images, places = prepare(results, profile, td)
        builders = BUILDERS
        if profile.get("template"):
            import tpl_builder                     # designed-page styles
            builders = tpl_builder.BUILDERS
        for kind in wanted_outputs(profile, asked):
            builder = builders.get(kind)
            if not builder:
                continue
            dest = OUT / f"{stem}.{kind}"
            try:
                builder(results, images, places, profile, title, dest)
            except ImportError as e:
                # PPTX needs python-pptx, which lives in the WEB layer's
                # requirements (requirements.txt is frozen). A bare CLI install
                # must still get its PDF, and must be told why it got only that
                # rather than watching the whole build die (rule 17).
                print(f"[report] cannot build .{kind}: {e} — install it with "
                      "`pip install -r requirements-web.txt`. The other "
                      "outputs were still written.", flush=True)
                dest.unlink(missing_ok=True)
                continue
            progress.wrote(dest, _mb(dest))
    extra = f"  ({skipped} skipped)" if skipped else ""
    print(f"[report] {len(results)} post(s) in {slug}{extra}", flush=True)


if __name__ == "__main__":
    main()
